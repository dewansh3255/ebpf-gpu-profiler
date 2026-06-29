#!/usr/bin/env python3
"""
G_21_make_presentation.py
Generates the Group 21 GRS Part B PowerPoint presentation.
Colors: IIITD teal #2AB5AE, dark gray #595959, white background.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Paths ──────────────────────────────────────────────────────────────────────
SCRIPT_DIR   = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH    = os.path.join(SCRIPT_DIR, "images.png")
PLOTS_DIR    = os.path.join(SCRIPT_DIR, "results", "plots_hardcoded")
OUT_PATH     = os.path.join(SCRIPT_DIR, "results", "G_21_Presentation.pptx")

# ── Palette ───────────────────────────────────────────────────────────────────
TEAL        = RGBColor(0x2A, 0xB5, 0xAE)
DARK_GRAY   = RGBColor(0x59, 0x59, 0x59)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)
LIGHT_TEAL  = RGBColor(0xD6, 0xF2, 0xF1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
ORANGE      = RGBColor(0xE8, 0x7A, 0x1E)

# ── Slide dimensions (widescreen 13.33 × 7.5 in) ─────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

BLANK_LAYOUT = prs.slide_layouts[6]   # truly blank


# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def add_logo(slide):
    """Place IIITD logo top-right on every slide."""
    if os.path.exists(LOGO_PATH):
        w, h = Inches(1.35), Inches(0.65)
        slide.shapes.add_picture(LOGO_PATH,
                                 SW - w - Inches(0.18),
                                 Inches(0.12), w, h)

def add_slide_number(slide, num):
    """Small page number bottom-right."""
    tb = slide.shapes.add_textbox(SW - Inches(0.7), SH - Inches(0.35),
                                  Inches(0.5), Inches(0.28))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(9)
    run.font.color.rgb = MID_GRAY

def add_top_bar(slide, color=TEAL, height=Inches(0.06)):
    """Thin coloured rule at very top of slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SW, height)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_bottom_bar(slide, color=TEAL, height=Inches(0.04)):
    bar = slide.shapes.add_shape(1, 0, SH - height, SW, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def txbox(slide, text, left, top, width, height,
          font_size=14, bold=False, color=DARK_GRAY,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def section_title(slide, text, top=Inches(0.72)):
    """Section heading with teal left accent bar."""
    bar = slide.shapes.add_shape(1, Inches(0.35), top,
                                  Inches(0.07), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    txbox(slide, text,
          left=Inches(0.52), top=top - Inches(0.02),
          width=Inches(10), height=Inches(0.55),
          font_size=22, bold=True, color=DARK_GRAY)

def bullet_block(slide, bullets, left, top, width, height,
                 font_size=13, indent=False, color=DARK_GRAY):
    """Add a list of bullet strings as separate text boxes (clean layout)."""
    line_h = Inches(0.36)
    for i, bullet in enumerate(bullets):
        prefix = "  •  " if not indent else "       –  "
        txbox(slide, prefix + bullet,
              left=left, top=top + i * line_h,
              width=width, height=line_h,
              font_size=font_size, color=color)

def add_image(slide, path, left, top, width, height=None):
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, left, top, width, height)
        else:
            slide.shapes.add_picture(path, left, top, width)

def teal_box(slide, left, top, width, height, text, font_size=12):
    """Teal filled box with white text — for metric cards."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = TEAL
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE

def gray_box(slide, left, top, width, height, text,
             font_size=12, text_color=DARK_GRAY, bg=LIGHT_TEAL):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = TEAL
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = False
    run.font.color.rgb = text_color

def metric_card(slide, left, top, w, h, label, value, unit=""):
    teal_box(slide, left, top, w, Inches(0.38), label, font_size=10)
    gray_box(slide, left, top + Inches(0.38), w, h - Inches(0.38),
             f"{value}\n{unit}", font_size=13, text_color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide, height=Inches(0.08))
    add_bottom_bar(slide)

    # Teal accent block left
    accent = slide.shapes.add_shape(1, 0, 0, Inches(0.55), SH)
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Logo
    add_logo(slide)

    # Institute name
    txbox(slide, "Indraprastha Institute of Information Technology Delhi",
          Inches(0.75), Inches(0.2), Inches(10.5), Inches(0.45),
          font_size=13, color=MID_GRAY, italic=True)

    # Main title
    txbox(slide, "Profiling CPU, Network Stack, and GPU Overheads\nin Containerized vs Non-Containerized ML Workloads\nusing eBPF and eGPU",
          Inches(0.75), Inches(1.3), Inches(11.2), Inches(2.2),
          font_size=30, bold=True, color=DARK_GRAY)

    # Horizontal rule
    rule = slide.shapes.add_shape(1, Inches(0.75), Inches(3.45),
                                   Inches(11.2), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()

    # Subtitle
    txbox(slide, "GRS Project Part B  |  Group 21",
          Inches(0.75), Inches(3.6), Inches(8), Inches(0.45),
          font_size=16, bold=True, color=TEAL)

    # Team members
    members = [
        "Dewansh Khandelwal", "Palak Mishra",
        "Sanskar Goyal",       "Yash Nimkar", "Kunal Verma"
    ]
    txbox(slide, "  •  " + "   •  ".join(members),
          Inches(0.75), Inches(4.15), Inches(11.5), Inches(0.5),
          font_size=13, color=DARK_GRAY)

    # Course info
    txbox(slide, "April 2026",
          Inches(0.75), Inches(4.75), Inches(5), Inches(0.4),
          font_size=13, color=MID_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation
# ══════════════════════════════════════════════════════════════════════════════
def slide_motivation(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Motivation")

    # Three motivation columns
    cols = [
        ("Containerisation Gap",
         ["Docker & Kubernetes are standard in ML clusters",
          "cgroup CPU limits add scheduling overhead",
          "veth bridge traversal on every packet",
          "No kernel-level breakdown by subsystem exists"]),
        ("TCP Network Blindspot",
         ["Most research clusters use RDMA, but small labs use TCP",
          "AllReduce over TCP — kernel-level cost unknown",
          "Socket buffer stalls invisible to ML framework",
          "eBPF can observe tcp_sendmsg without modifying PyTorch"]),
        ("eBPF Opportunity",
         ["Production-grade kernel observability",
          "kprobes, tracepoints, TC hooks, XDP, uprobes",
          "Near-zero overhead vs strace/ptrace",
          "No prior work unifies CPU + GPU + network timelines"]),
    ]

    col_w = Inches(3.9)
    for i, (title, bullets) in enumerate(cols):
        left = Inches(0.35) + i * (col_w + Inches(0.12))
        top  = Inches(1.45)
        # Header box
        teal_box(slide, left, top, col_w, Inches(0.48), title, font_size=13)
        # Body box
        body_h = Inches(4.5)
        body = slide.shapes.add_shape(1, left, top + Inches(0.48),
                                       col_w, body_h)
        body.fill.solid()
        body.fill.fore_color.rgb = LIGHT_TEAL
        body.line.color.rgb = TEAL
        for j, b in enumerate(bullets):
            txbox(slide, "•  " + b,
                  left + Inches(0.1),
                  top + Inches(0.6) + j * Inches(0.72),
                  col_w - Inches(0.2), Inches(0.7),
                  font_size=12, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Objectives
# ══════════════════════════════════════════════════════════════════════════════
def slide_objectives(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Objectives")

    objs = [
        ("01", "Containerization Overhead Quantification",
         "Measure CPU scheduling delays, syscall counts/latency, and network stack traversal between native and Docker execution of a 2-GPU ML job.\nInstrument: sched_switch, sys_enter/exit, tcp_sendmsg/recvmsg, nvidia-smi.\nDeliverable: Per-subsystem overhead table."),
        ("02", "TCP Network Stack Profiling (Two-Node)",
         "Deploy PyTorch FL across two laptops over TCP and profile gradient sync cost at the socket layer.\nInstrument: tcp_sendmsg, tcp_recvmsg, tcp_retransmit_skb, sched_wakeup.\nDeliverable: Timeline showing GPU compute, TCP send/recv, and scheduler events."),
        ("03", "Cross-Scenario Comparison & Overhead Attribution",
         "Compare: (a) single-server bare-metal, (b) single-server containerized, (c) distributed TCP.\nFor each: GPU util %, AllReduce latency, scheduler delay, network overhead, syscall frequency.\nDeliverable: Comparative profiles + conclusion on where to optimize."),
    ]

    for i, (num_lbl, title, desc) in enumerate(objs):
        top = Inches(1.45) + i * Inches(1.82)
        # Number badge
        badge = slide.shapes.add_shape(1, Inches(0.35), top,
                                        Inches(0.55), Inches(1.3))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        txbox(slide, num_lbl,
              Inches(0.35), top + Inches(0.35),
              Inches(0.55), Inches(0.55),
              font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        txbox(slide, title,
              Inches(1.05), top + Inches(0.05),
              Inches(11.5), Inches(0.45),
              font_size=14, bold=True, color=TEAL)
        # Desc
        txbox(slide, desc,
              Inches(1.05), top + Inches(0.48),
              Inches(11.5), Inches(0.9),
              font_size=11.5, color=DARK_GRAY)
        # Separator line
        if i < 2:
            rule = slide.shapes.add_shape(1, Inches(0.35),
                                           top + Inches(1.65),
                                           Inches(12.6), Inches(0.02))
            rule.fill.solid()
            rule.fill.fore_color.rgb = LIGHT_TEAL
            rule.line.fill.background()

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Background
# ══════════════════════════════════════════════════════════════════════════════
def slide_background(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Background")

    panels = [
        ("eBPF",
         ["Sandboxed kernel programs — no module needed",
          "Attachment: kprobes, tracepoints, uprobes, TC hooks, XDP",
          "Maps (hash, ringbuf) for kernel→user data transfer",
          "Frontend: BCC (Python), bpftrace, libbpf",
          "Near-zero overhead vs strace/ptrace"]),
        ("eGPU",
         ["First eBPF runtime that offloads to GPU via PTX injection",
          "Instruments kernel launches, LD/ST memory ops",
          "Events in shared eBPF maps — no copy overhead",
          "Requires LLVM ≥15 + bpftime runtime",
          "We use hybrid fallback: uprobes on libcuda.so"]),
        ("Containerisation",
         ["Namespaces: network, pid, mount isolation",
          "cgroups: per-cgroup CPU budgets → throttle events",
          "Docker bridge: container netns → veth → docker0 → host",
          "≥2 namespace crossings per packet vs bare-metal",
          "GPU pass-through via NVIDIA Container Toolkit"]),
        ("PyTorch DDP",
         ["Data Parallelism: each GPU holds full model replica",
          "AllReduce aggregates gradients after each backward pass",
          "NCCL: PCIe on single-server",
          "Gloo/TCP: cross-node gradient sync over sockets",
          "No PyTorch/NCCL source modification needed"]),
    ]

    col_w = Inches(3.05)
    for i, (title, bullets) in enumerate(panels):
        left = Inches(0.3) + i * (col_w + Inches(0.12))
        top  = Inches(1.42)
        teal_box(slide, left, top, col_w, Inches(0.44), title, font_size=13)
        body = slide.shapes.add_shape(1, left, top + Inches(0.44),
                                       col_w, Inches(4.6))
        body.fill.solid()
        body.fill.fore_color.rgb = LIGHT_TEAL
        body.line.color.rgb = TEAL
        for j, b in enumerate(bullets):
            txbox(slide, "•  " + b,
                  left + Inches(0.1),
                  top + Inches(0.55) + j * Inches(0.82),
                  col_w - Inches(0.18), Inches(0.8),
                  font_size=11, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Design & Implementation (overview)
# ══════════════════════════════════════════════════════════════════════════════
def slide_design_overview(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Design & Implementation — System Overview")

    # Hardware configs
    txbox(slide, "Experiment Configurations",
          Inches(0.35), Inches(1.42), Inches(12.5), Inches(0.38),
          font_size=14, bold=True, color=DARK_GRAY)

    cfg_data = [
        ("Config A\nSingle-Server\nMulti-GPU",
         "1 server × 2 NVIDIA H100 NVL\nAMD EPYC 9354 32-core\n503 GB RAM, CUDA 12.1\nKernel 5.15 / Ubuntu 22.04"),
        ("Config B\nTwo-Node\nDistributed",
         "2 laptops × 1 RTX GPU each\nConnected over Wi-Fi / TCP\nFL workload (HTTP parameter server)\nPyTorch ResNet-18 / CIFAR-10"),
    ]

    for i, (label, detail) in enumerate(cfg_data):
        left = Inches(0.35) + i * Inches(6.3)
        teal_box(slide, left, Inches(1.85), Inches(2.2), Inches(0.9), label, font_size=12)
        gray_box(slide, left + Inches(2.25), Inches(1.85), Inches(3.85), Inches(0.9),
                 detail, font_size=11)

    # Four profiling components
    txbox(slide, "Profiling Framework — 4 Concurrent Layers",
          Inches(0.35), Inches(3.0), Inches(12.5), Inches(0.38),
          font_size=14, bold=True, color=DARK_GRAY)

    comps = [
        ("ML Workload", "ResNet-18 / CIFAR-10\nPyTorch DDP, 2 GPU\n10 epochs, batch=128"),
        ("eBPF Probing", "sched_switch tracepoint\nsys_enter/exit tracepoints\ntcp_sendmsg/recvmsg kprobes"),
        ("GPU Profiling", "Primary: eGPU uprobes\non libcuda.so (hybrid)\nSecondary: nvidia-smi 100ms"),
        ("Data Collection", "Ring buffer → userspace\nTimestamped CSV/JSON\nCLOCK_MONOTONIC aligned"),
    ]

    comp_w = Inches(3.05)
    for i, (title, detail) in enumerate(comps):
        left = Inches(0.3) + i * (comp_w + Inches(0.12))
        teal_box(slide, left, Inches(3.45), comp_w, Inches(0.44), title, font_size=12)
        gray_box(slide, left, Inches(3.89), comp_w, Inches(1.1), detail, font_size=11)

    # Workload detail
    txbox(slide, "Workload: ResNet-18 on CIFAR-10  |  PyTorch DDP  |  torchrun --nproc_per_node=2  |  10 epochs  |  batch_size=128  |  lr=0.01",
          Inches(0.35), Inches(5.2), Inches(12.5), Inches(0.4),
          font_size=11, color=MID_GRAY, italic=True)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Implementation: eBPF Profilers
# ══════════════════════════════════════════════════════════════════════════════
def slide_impl_ebpf(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Implementation — eBPF Profiling Layer")

    profilers = [
        ("CPU Scheduler\nG_21_cpu_profiler.py",
         ["Attaches to sched:sched_switch tracepoint",
          "Records: scheduling latency (run-queue wait time)",
          "Context switch count per CPU",
          "Per-process execution runtime",
          "Events streamed via perf buffers → CSV"]),
        ("Syscall Counter\nG_21_syscall_counter.py",
         ["Attaches to raw_syscalls:sys_enter / sys_exit",
          "Per-syscall: count, avg/min/max latency",
          "Key calls: futex, gettid, poll, read, ioctl",
          "Reveals namespace-crossing overhead",
          "Outputs: CSV with 121–167 unique syscall types"]),
        ("Network Profiler\nG_21_net_profiler.py",
         ["kprobes on tcp_sendmsg + tcp_recvmsg",
          "Measures: start→end latency per TCP operation",
          "Captures: bytes, PID, comm, direction",
          "Shows container veth bridge overhead",
          "Outputs: per-event CSV with latency in μs"]),
        ("GPU Monitor\nG_21_gpu_monitor_nvidia.py",
         ["nvidia-smi polling at 100ms intervals",
          "Metrics: util%, memory MB, power W, temp °C",
          "Per-GPU (both GPU 0 and GPU 1 captured)",
          "eGPU hybrid: uprobes on libcuda.so",
          "Traces cuLaunchKernel, cuMemcpyHtoD/DtoH"]),
    ]

    col_w = Inches(3.05)
    for i, (title, bullets) in enumerate(profilers):
        left = Inches(0.3) + i * (col_w + Inches(0.12))
        top  = Inches(1.42)
        teal_box(slide, left, top, col_w, Inches(0.6), title, font_size=11)
        body = slide.shapes.add_shape(1, left, top + Inches(0.6),
                                       col_w, Inches(4.4))
        body.fill.solid()
        body.fill.fore_color.rgb = LIGHT_TEAL
        body.line.color.rgb = TEAL
        for j, b in enumerate(bullets):
            txbox(slide, "•  " + b,
                  left + Inches(0.1),
                  top + Inches(0.72) + j * Inches(0.78),
                  col_w - Inches(0.18), Inches(0.76),
                  font_size=11, color=DARK_GRAY)

    # Key note about BCC Python path
    txbox(slide, "⚙  BCC eBPF scripts use system Python with PYTHONPATH=/usr/lib/python3/dist-packages (not venv) — venv lacks BCC bindings",
          Inches(0.35), Inches(6.45), Inches(12.5), Inches(0.38),
          font_size=10.5, color=MID_GRAY, italic=True)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — eGPU: Attempted + Hybrid Fallback
# ══════════════════════════════════════════════════════════════════════════════
def slide_egpu(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "eGPU: Attempted Instrumentation & Hybrid Fallback")

    # Left: Intended eGPU
    teal_box(slide, Inches(0.35), Inches(1.42), Inches(6.05), Inches(0.48),
             "eGPU (bpftime) — Intended Approach", font_size=13)
    intended = [
        "PTX injection into running GPU kernels",
        "Instruments LD/ST memory operations",
        "Kernel launch/exit boundary tracing",
        "Events in shared eBPF maps (zero copy)",
        "Requires LLVM ≥15 + bpftime runtime + Frida",
    ]
    body_l = slide.shapes.add_shape(1, Inches(0.35), Inches(1.9),
                                     Inches(6.05), Inches(2.2))
    body_l.fill.solid()
    body_l.fill.fore_color.rgb = LIGHT_TEAL
    body_l.line.color.rgb = TEAL
    for j, b in enumerate(intended):
        txbox(slide, "•  " + b, Inches(0.45), Inches(2.02) + j * Inches(0.4),
              Inches(5.8), Inches(0.38), font_size=12, color=DARK_GRAY)

    # Blockers box
    box = slide.shapes.add_shape(1, Inches(0.35), Inches(4.2),
                                  Inches(6.05), Inches(1.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE8)
    box.line.color.rgb = ORANGE
    txbox(slide, "⚠  Build Blockers",
          Inches(0.45), Inches(4.25), Inches(5.9), Inches(0.38),
          font_size=13, bold=True, color=ORANGE)
    blockers = [
        "Host LLVM 14 (requires ≥15); Docker image has LLVM 10",
        "PyTorch/cuDNN kernels ship as cubins — no PTX available",
        "bpftime aborts when PTX is absent",
        "nsys failed to attach to torchrun child CUDA contexts",
    ]
    for j, b in enumerate(blockers):
        txbox(slide, "✗  " + b, Inches(0.45), Inches(4.68) + j * Inches(0.36),
              Inches(5.9), Inches(0.35), font_size=11.5,
              color=RGBColor(0xC0, 0x40, 0x00))

    # Right: Hybrid approach
    teal_box(slide, Inches(6.6), Inches(1.42), Inches(6.35), Inches(0.48),
             "Hybrid Fallback — What We Did Instead", font_size=13)
    hybrid = [
        "eBPF uprobes on libcuda.so (CPU-side tracing)",
        "Instruments: cuLaunchKernel, cuMemcpyHtoD_v2, cuMemcpyDtoH_v2",
        "Captures: kernel launch timestamps, memory transfer durations",
        "Events tagged COMPUTE_MATH or MEM_TRANSFER",
        "nvidia-smi polling: util%, power W, temp °C at 100ms",
        "CLOCK_MONOTONIC timestamps → aligned with eBPF events",
    ]
    body_r = slide.shapes.add_shape(1, Inches(6.6), Inches(1.9),
                                     Inches(6.35), Inches(2.6))
    body_r.fill.solid()
    body_r.fill.fore_color.rgb = LIGHT_TEAL
    body_r.line.color.rgb = TEAL
    for j, b in enumerate(hybrid):
        txbox(slide, "✓  " + b, Inches(6.7), Inches(2.02) + j * Inches(0.42),
              Inches(6.1), Inches(0.4), font_size=12, color=DARK_GRAY)

    txbox(slide, "Limitation: Traces CUDA driver API calls from CPU side — not PTX-level GPU kernel instrumentation. GPU execution time is inferred, not directly measured.",
          Inches(6.6), Inches(4.62), Inches(6.35), Inches(0.7),
          font_size=11, italic=True, color=MID_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Evaluation: Training Results
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_training(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — Training Results (ResNet-18 / CIFAR-10)")

    # Key numbers row
    cards = [
        ("Native — Total Time", "31.8 s"),
        ("Container — Total Time", "34.3 s"),
        ("Time Overhead", "+7.9%"),
        ("Final Accuracy (both)", "81.1%"),
        ("Native Throughput", "9,274 samp/s"),
        ("Container Throughput", "8,472 samp/s"),
    ]
    cw = Inches(2.08)
    for i, (lbl, val) in enumerate(cards):
        metric_card(slide, Inches(0.35) + i * (cw + Inches(0.05)),
                    Inches(1.42), cw, Inches(0.9), lbl, val)

    # Two plots side by side
    add_image(slide,
              os.path.join(PLOTS_DIR, "01_training_loss.png"),
              Inches(0.35), Inches(2.48), Inches(6.2))
    add_image(slide,
              os.path.join(PLOTS_DIR, "02_training_accuracy.png"),
              Inches(6.75), Inches(2.48), Inches(6.2))

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Evaluation: Throughput & Epoch Time
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_throughput(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — Throughput & Epoch Time")

    add_image(slide,
              os.path.join(PLOTS_DIR, "03_throughput.png"),
              Inches(0.35), Inches(1.38), Inches(6.3))
    add_image(slide,
              os.path.join(PLOTS_DIR, "04_epoch_time.png"),
              Inches(6.75), Inches(1.38), Inches(6.2))

    txbox(slide,
          "Native epochs average 2.65s each; Container averages 2.98s (+12.5% per epoch). Epoch 1 is slowest in both (CUDA init). Container variance is lower — cgroup scheduling is more deterministic.",
          Inches(0.35), Inches(6.42), Inches(12.5), Inches(0.42),
          font_size=11.5, italic=True, color=MID_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Evaluation: Total Comparison
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_total(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — Overall Performance Comparison")

    add_image(slide,
              os.path.join(PLOTS_DIR, "05_total_comparison.png"),
              Inches(0.35), Inches(1.38), Inches(8.5))

    # Side callouts
    callouts = [
        ("+7.9%", "Training time overhead"),
        ("-8.6%", "Throughput reduction"),
        ("-3.4%", "GPU util (active %)"),
        ("+25.6%", "Power consumption ⚡"),
    ]
    for i, (val, label) in enumerate(callouts):
        top = Inches(1.6) + i * Inches(1.22)
        color = TEAL if not val.startswith("+2") else ORANGE
        teal_box(slide, Inches(9.05), top, Inches(1.6), Inches(0.42), val,
                 font_size=15)
        txbox(slide, label, Inches(9.05), top + Inches(0.44),
              Inches(3.85), Inches(0.42), font_size=11, color=DARK_GRAY)

    txbox(slide,
          "GPU remains the bottleneck in both modes — utilization drop of only 3.4%.\nThe notable overhead is power (+25.6%): container startup and namespace overhead keeps GPU power draw elevated even during idle periods.",
          Inches(9.05), Inches(6.45), Inches(4.0), Inches(0.85),
          font_size=11, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Evaluation: GPU Metrics
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_gpu(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — GPU Utilization, Power & Temperature")

    add_image(slide,
              os.path.join(PLOTS_DIR, "06_gpu_metrics.png"),
              Inches(0.35), Inches(1.38), Inches(12.6))

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Evaluation: Syscall Analysis
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_syscall(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — Syscall Analysis")

    add_image(slide,
              os.path.join(PLOTS_DIR, "07_syscall_comparison.png"),
              Inches(0.35), Inches(1.38), Inches(6.3))
    add_image(slide,
              os.path.join(PLOTS_DIR, "08_syscall_latency.png"),
              Inches(6.75), Inches(1.38), Inches(6.2))

    cards = [
        ("Native — Total Syscalls", "11,926,534"),
        ("Container — Total Syscalls", "8,180,511"),
        ("Native — Unique Types", "121"),
        ("Container — Unique Types", "167  (+38%)"),
    ]
    cw = Inches(3.1)
    for i, (lbl, val) in enumerate(cards):
        metric_card(slide, Inches(0.35) + i * (cw + Inches(0.07)),
                    Inches(6.28), cw, Inches(0.88), lbl, val, "")

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Evaluation: Network Analysis
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_network(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — Network Stack Analysis")

    add_image(slide,
              os.path.join(PLOTS_DIR, "09_network_comparison.png"),
              Inches(0.35), Inches(1.38), Inches(7.5))

    # Right side: distributed TCP data
    teal_box(slide, Inches(8.1), Inches(1.42), Inches(4.85), Inches(0.44),
             "Distributed TCP (Sanskar-Kunal)", font_size=12)
    dist_stats = [
        ("TCP Send events",  "575",    "(avg 21.5 MB each)"),
        ("TCP Recv events",  "69,764", "(HTTP polling 0.25 MB avg)"),
        ("TCP Retransmits",  "292",    "(Wi-Fi packet drops)"),
        ("Sched delay mean", "2,250 μs", "(>1ms filtered)"),
        ("Sched delay max",  "22,058 μs",""),
        ("Run duration",     "~8 min 53 s",""),
    ]
    body = slide.shapes.add_shape(1, Inches(8.1), Inches(1.86),
                                   Inches(4.85), Inches(3.4))
    body.fill.solid()
    body.fill.fore_color.rgb = LIGHT_TEAL
    body.line.color.rgb = TEAL
    for j, (lbl, val, note) in enumerate(dist_stats):
        txbox(slide, f"{lbl}:", Inches(8.2), Inches(2.0) + j * Inches(0.52),
              Inches(2.1), Inches(0.48), font_size=11, bold=True, color=DARK_GRAY)
        txbox(slide, f"{val}  {note}", Inches(10.35), Inches(2.0) + j * Inches(0.52),
              Inches(2.5), Inches(0.48), font_size=11, color=TEAL)

    txbox(slide,
          "292 retransmits = Wi-Fi unreliability causing gradient sync stalls.\n"
          "Native TCP send avg 30.4μs vs Container 35.7μs (+17.4%).",
          Inches(0.35), Inches(6.42), Inches(12.5), Inches(0.55),
          font_size=11.5, italic=True, color=MID_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Evaluation: Scheduler Latency & Context Switches
# ══════════════════════════════════════════════════════════════════════════════
def slide_eval_sched(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Evaluation — CPU Scheduler Latency & Context Switches")

    add_image(slide,
              os.path.join(PLOTS_DIR, "10_sched_latency.png"),
              Inches(0.35), Inches(1.38), Inches(8.6))

    cards = [
        ("Native Mean Latency",    "13.9 μs"),
        ("Container Mean Latency", "17.7 μs"),
        ("Distributed Mean",       "2,250 μs"),
        ("Native Ctx Switches",    "3,367,166"),
        ("Container Ctx Switches", "2,508,117"),
        ("Native Rate",            "55,108 /s"),
    ]
    cw = Inches(2.0)
    for i, (lbl, val) in enumerate(cards[:3]):
        metric_card(slide, Inches(9.1) + i * (cw + Inches(0.06)),
                    Inches(1.5), cw, Inches(0.95), lbl, val)

    for i, (lbl, val) in enumerate(cards[3:]):
        metric_card(slide, Inches(9.1) + i * (cw + Inches(0.06)),
                    Inches(2.55), cw, Inches(0.95), lbl, val)

    txbox(slide,
          "Container scheduling latency mean is +27.3% higher (17.7 vs 13.9 μs),\n"
          "but P95 is actually lower (14.4 vs 15.2 μs) — cgroup scheduling is more uniform.\n"
          "Distributed TCP sched latency is 160× higher due to I/O blocking on socket waits.",
          Inches(9.1), Inches(3.65), Inches(3.95), Inches(1.1),
          font_size=11, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Overhead Summary Table
# ══════════════════════════════════════════════════════════════════════════════
def slide_summary_table(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Overhead Summary — All Metrics")

    add_image(slide,
              os.path.join(PLOTS_DIR, "11_overhead_summary_table.png"),
              Inches(0.35), Inches(1.38), Inches(12.6))

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Insights & Findings
# ══════════════════════════════════════════════════════════════════════════════
def slide_insights(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Insights & Key Findings")

    findings = [
        ("Container adds ~7.9% training time overhead",
         "On H100 NVL GPUs, Docker adds 2.5 extra seconds per 31.8s run. GPU is the bottleneck — not container overhead."),
        ("GPU utilization is near-identical (76.7% vs 74.1%)",
         "Active-period GPU util drops only 3.4%. Both environments are GPU-bound. Containerization doesn't throttle the GPU."),
        ("Power consumption spikes +25.6% in container",
         "Container init (overlay FS, namespace setup) keeps GPU power elevated during idle phases (198.9W vs 158.4W avg)."),
        ("Syscall diversity increases +38% in containers",
         "167 unique syscall types in container vs 121 native. Overlay FS, veth, and cgroup management generate extra kernel calls."),
        ("Scheduling latency mean +27.3% in container",
         "17.7 vs 13.9 μs mean. But P95 is slightly better in container (14.4 vs 15.2 μs) — cgroup scheduling is more uniform."),
        ("Distributed TCP sched latency is 160× worse",
         "2,250 μs mean in distributed FL setup (vs 13.9 μs native). Wi-Fi blocking I/O completely dominates CPU scheduling cost."),
        ("292 TCP retransmits in distributed setup",
         "Wi-Fi instability causes gradient sync stalls. RDMA/InfiniBand would eliminate this entirely."),
        ("eBPF is a viable unified profiling layer",
         "Single framework captures CPU scheduling, syscalls, TCP stack, and GPU driver events with near-zero overhead."),
    ]

    col1 = findings[:4]
    col2 = findings[4:]
    finding_h = Inches(1.32)
    col_w = Inches(6.2)

    for col_i, col in enumerate([col1, col2]):
        left = Inches(0.35) + col_i * (col_w + Inches(0.22))
        for j, (title, body) in enumerate(col):
            top = Inches(1.42) + j * finding_h
            # Number badge
            badge = slide.shapes.add_shape(1, left, top + Inches(0.04),
                                            Inches(0.38), Inches(0.38))
            badge.fill.solid()
            badge.fill.fore_color.rgb = TEAL
            badge.line.fill.background()
            txbox(slide, str(col_i * 4 + j + 1),
                  left, top + Inches(0.04),
                  Inches(0.38), Inches(0.38),
                  font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txbox(slide, title,
                  left + Inches(0.44), top + Inches(0.04),
                  col_w - Inches(0.44), Inches(0.38),
                  font_size=12, bold=True, color=TEAL)
            txbox(slide, body,
                  left + Inches(0.44), top + Inches(0.44),
                  col_w - Inches(0.44), Inches(0.82),
                  font_size=10.5, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Limitations & Future Work
# ══════════════════════════════════════════════════════════════════════════════
def slide_limitations(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "Limitations & Future Work")

    lims = [
        ("eGPU Build Failed",
         "PTX-level GPU instrumentation required LLVM ≥15. We used a CPU-side uprobe fallback. True GPU kernel instrumentation remains unachieved."),
        ("Single Trial",
         "Only one run per configuration (native/container). Statistical variance not captured. 3+ trials with median selection would strengthen conclusions."),
        ("eBPF Requires Root",
         "All eBPF profilers need sudo. Deployment in locked-down cloud environments (e.g., GKE) is not possible without privileged containers."),
        ("nvidia-smi Resolution",
         "100ms polling misses sub-millisecond GPU events. CUPTI or Nsight Compute would provide microsecond-level GPU kernel timing."),
        ("Single-Node Only",
         "Config A experiments run on one server. Multi-node DDP overhead (NCCL ring AllReduce across nodes) not measured."),
        ("System-Wide eBPF Probes",
         "Our eBPF profilers capture all processes, not just the training job. This inflates syscall and context-switch counts."),
    ]

    future = [
        "Upgrade to LLVM ≥15 and build full eGPU (bpftime) for PTX-level GPU tracing",
        "Run 5+ trials and report mean ± std for all metrics",
        "Add CUDA stream-level profiling via CUPTI for compute/communication overlap measurement",
        "Extend to multi-node DDP on InfiniBand vs TCP and quantify RDMA benefit",
        "Add per-PID eBPF filtering to isolate training process from system noise",
        "Test on smaller GPUs (RTX class) where container overhead may be relatively larger",
    ]

    # Limitations — left 2/3
    teal_box(slide, Inches(0.35), Inches(1.42), Inches(8.3), Inches(0.44),
             "Current Limitations", font_size=13)
    lim_h = Inches(0.88)
    for j, (title, body) in enumerate(lims):
        top = Inches(1.92) + j * lim_h
        txbox(slide, "⚠  " + title,
              Inches(0.45), top, Inches(8.1), Inches(0.32),
              font_size=12, bold=True, color=ORANGE)
        txbox(slide, body,
              Inches(0.45), top + Inches(0.32), Inches(8.1), Inches(0.52),
              font_size=11, color=DARK_GRAY)

    # Future Work — right 1/3
    teal_box(slide, Inches(8.85), Inches(1.42), Inches(4.1), Inches(0.44),
             "Future Work", font_size=13)
    body_r = slide.shapes.add_shape(1, Inches(8.85), Inches(1.86),
                                     Inches(4.1), Inches(5.0))
    body_r.fill.solid()
    body_r.fill.fore_color.rgb = LIGHT_TEAL
    body_r.line.color.rgb = TEAL
    for j, b in enumerate(future):
        txbox(slide, "→  " + b,
              Inches(8.95), Inches(2.0) + j * Inches(0.8),
              Inches(3.9), Inches(0.75),
              font_size=11, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
def slide_conclusion(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide, height=Inches(0.08))
    add_bottom_bar(slide)

    # Teal accent block left
    accent = slide.shapes.add_shape(1, 0, 0, Inches(0.55), SH)
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_logo(slide)

    txbox(slide, "Conclusions",
          Inches(0.75), Inches(0.55), Inches(10), Inches(0.55),
          font_size=28, bold=True, color=DARK_GRAY)

    rule = slide.shapes.add_shape(1, Inches(0.75), Inches(1.1),
                                   Inches(11.5), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()

    conclusions = [
        ("Containerization overhead is real but modest on H100s",
         "Docker adds 7.9% training time and -8.6% throughput. GPU utilization is almost identical — container overhead does not significantly throttle GPU-bound workloads."),
        ("Power is the hidden container cost",
         "Container GPU power is +25.6% higher. Container initialisation keeps the GPU at elevated power even during idle phases between epochs."),
        ("Syscall diversity is a container fingerprint",
         "+38% unique syscall types in containers reveal overlay FS, veth, and cgroup activity invisible at the application layer."),
        ("Distributed TCP over Wi-Fi is the dominant bottleneck",
         "160× higher scheduling latency and 292 retransmits vs zero in native/container. Network infrastructure quality matters far more than containerization."),
        ("eBPF is a unified, zero-overhead profiling layer",
         "A single eBPF framework successfully profiles CPU scheduling, syscalls, TCP stack, and GPU driver events concurrently — validating eBPF as the right tool for ML system observability."),
    ]

    for j, (title, body) in enumerate(conclusions):
        top = Inches(1.25) + j * Inches(1.14)
        badge = slide.shapes.add_shape(1, Inches(0.75), top + Inches(0.04),
                                        Inches(0.38), Inches(0.38))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        txbox(slide, str(j + 1),
              Inches(0.75), top + Inches(0.04), Inches(0.38), Inches(0.38),
              font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, title,
              Inches(1.25), top + Inches(0.04), Inches(11.5), Inches(0.38),
              font_size=13, bold=True, color=TEAL)
        txbox(slide, body,
              Inches(1.25), top + Inches(0.44), Inches(11.5), Inches(0.62),
              font_size=11.5, color=DARK_GRAY)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — References & Thank You
# ══════════════════════════════════════════════════════════════════════════════
def slide_thanks(num):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_bar(slide)
    add_bottom_bar(slide)
    add_logo(slide)
    section_title(slide, "References")

    refs = [
        "[1] Yang et al., \"eGPU: Extending eBPF Programmability and Observability to GPUs,\" HCDS '25. https://dl.acm.org/doi/pdf/10.1145/3723851.3726984",
        "[2] eGPU (bpftime) GitHub: https://github.com/eunomia-bpf/eGPU",
        "[3] BCC (BPF Compiler Collection): https://github.com/iovisor/bcc",
        "[4] NVIDIA Nsight Systems: https://developer.nvidia.com/nsight-systems",
        "[5] CUDA CUPTI: https://docs.nvidia.com/cupti/",
        "[6] PyTorch ResNet: https://docs.pytorch.org/vision/main/models/resnet.html",
        "[7] NVIDIA Container Toolkit: https://docs.nvidia.com/datacenter/cloud-native/container-toolkit/",
        "[8] PyTorch Distributed Data Parallel: https://pytorch.org/tutorials/intermediate/ddp_tutorial.html",
    ]
    for j, r in enumerate(refs):
        txbox(slide, r, Inches(0.45), Inches(1.45) + j * Inches(0.52),
              Inches(12.5), Inches(0.5), font_size=11, color=DARK_GRAY)

    # Thank you box
    ty = slide.shapes.add_shape(1, Inches(3.5), Inches(5.7), Inches(6.33), Inches(1.3))
    ty.fill.solid()
    ty.fill.fore_color.rgb = TEAL
    ty.line.fill.background()
    txbox(slide, "Thank You",
          Inches(3.5), Inches(5.78), Inches(6.33), Inches(0.65),
          font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, "Group 21  •  GRS Part B  •  IIIT Delhi",
          Inches(3.5), Inches(6.42), Inches(6.33), Inches(0.4),
          font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_slide_number(slide, num)


# ══════════════════════════════════════════════════════════════════════════════
# Build all slides
# ══════════════════════════════════════════════════════════════════════════════
slide_title(1)
slide_motivation(2)
slide_objectives(3)
slide_background(4)
slide_design_overview(5)
slide_impl_ebpf(6)
slide_egpu(7)
slide_eval_training(8)
slide_eval_throughput(9)
slide_eval_total(10)
slide_eval_gpu(11)
slide_eval_syscall(12)
slide_eval_network(13)
slide_eval_sched(14)
slide_summary_table(15)
slide_insights(16)
slide_limitations(17)
slide_conclusion(18)
slide_thanks(19)

os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"✓  Saved {len(prs.slides)} slides → {OUT_PATH}")
